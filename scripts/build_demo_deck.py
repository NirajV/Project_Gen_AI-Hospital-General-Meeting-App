"""
Generate the BioMedMeet demo-request PowerPoint deck.

Output: /app/docs/BioMedMeet_Demo_Deck.pptx (10 slides, 16:9) by default.

Personalisation (recommended):
    python scripts/build_demo_deck.py \
        --hospital-name "Ochsner Health" \
        --month May --year 2026

The hospital name appears on the cover ("Prepared for Ochsner Health · May 2026"),
the overview hook, the pricing card title, and the closing CTA. The output
filename auto-derives to `BioMedMeet_Demo_Deck_Ochsner_Health.pptx`.

Design system mirrors the marketing site:
  navy  #0b0b30  – primary text + dark CTAs
  green #3b6658  – feature accents
  olive #694e20  – pricing / hospital tone
  purple#68517d  – participants / data accents
  ivory #f9f5ee  – warm slide background
  white #ffffff
"""
import argparse
import re
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ----------------------------------------------------------------------------- 
# Brand palette + helpers
# ----------------------------------------------------------------------------- 

NAVY   = RGBColor(0x0B, 0x0B, 0x30)
GREEN  = RGBColor(0x3B, 0x66, 0x58)
OLIVE  = RGBColor(0x69, 0x4E, 0x20)
PURPLE = RGBColor(0x68, 0x51, 0x7D)
IVORY  = RGBColor(0xF9, 0xF5, 0xEE)
CREAM  = RGBColor(0xFD, 0xFA, 0xF3)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SLATE  = RGBColor(0x4A, 0x55, 0x68)
LIGHT  = RGBColor(0xEC, 0xEA, 0xE3)

# Tinted backgrounds matching the site card system
TINT_GREEN  = RGBColor(0xE8, 0xF5, 0xF0)
TINT_OLIVE  = RGBColor(0xF5, 0xF0, 0xE8)
TINT_PURPLE = RGBColor(0xF3, 0xED, 0xF5)
TINT_NAVY   = RGBColor(0xE8, 0xE8, 0xF5)

OUT_PATH = Path("/app/docs/BioMedMeet_Demo_Deck.pptx")
ASSETS_DIR = Path(__file__).resolve().parent / "deck_assets"
ASSETS_DIR.mkdir(exist_ok=True)


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    """Add a filled rectangle. Returns the shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        if line_width is not None:
            shape.line.width = line_width
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.12
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             size=14, color=NAVY, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri"):
    """Add a text box and return it."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def draw_logo(slide, x, y, *, height=Inches(0.5), light=False):
    """Draw the heartbeat/activity icon + 'BioMedMeet' wordmark."""
    h = height
    # Heart-line icon — simulated with 5 connected line segments (free-form).
    line_color = WHITE if light else NAVY
    box = h
    # Background pill for the icon
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box, box)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = WHITE if light else IVORY
    pill.line.color.rgb = line_color
    pill.line.width = Pt(1.5)
    pill.shadow.inherit = False

    # Draw a small heartbeat polyline as an actual freeform shape
    freeform = slide.shapes.build_freeform(int(x + box * 0.18), int(y + box * 0.50))
    seg_w = (box * 0.64) / 6
    pts = [
        (seg_w * 1, 0),
        (seg_w * 2, -box * 0.18),
        (seg_w * 3, box * 0.22),
        (seg_w * 4, -box * 0.30),
        (seg_w * 5, 0),
        (seg_w * 6, 0),
    ]
    freeform.add_line_segments(pts, close=False)
    line_shape = freeform.convert_to_shape()
    line_shape.fill.background()
    line_shape.line.color.rgb = line_color
    line_shape.line.width = Pt(2.5)

    # Wordmark
    text_color = WHITE if light else NAVY
    sub_color = LIGHT if light else SLATE
    tb = slide.shapes.add_textbox(x + box + Inches(0.12), y, Inches(3.5), h)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "BioMedMeet"
    r.font.name = "Calibri"
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = text_color


def add_footer(slide, page_num, total):
    """Bottom-right page number + small 'BioMedMeet · biomedmeet.com'."""
    add_text(slide, Inches(0.5), Inches(7.05), Inches(6), Inches(0.3),
             "BioMedMeet  ·  biomedmeet.com  ·  Demo@BioMedMeet.com",
             size=9, color=SLATE)
    add_text(slide, Inches(11.2), Inches(7.05), Inches(2), Inches(0.3),
             f"{page_num} / {total}", size=9, color=SLATE, align=PP_ALIGN.RIGHT)


# ----------------------------------------------------------------------------- 
# Slide builders
# ----------------------------------------------------------------------------- 

def make_cover(prs, *, hospital_name=None, prepared_for_line=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Deep-navy background
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY)

    # Decorative diagonal band (olive)
    band = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                  Inches(9.2), Inches(-1.5),
                                  Inches(8), Inches(11.0))
    band.rotation = 12
    band.fill.solid()
    band.fill.fore_color.rgb = OLIVE
    band.line.fill.background()
    band.shadow.inherit = False

    # Subtle accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                    Inches(7.6), Inches(-2.5),
                                    Inches(0.5), Inches(13))
    stripe.rotation = 12
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = GREEN
    stripe.line.fill.background()
    stripe.shadow.inherit = False

    # Logo (light variant) at the top-left
    draw_logo(slide, Inches(0.6), Inches(0.5), height=Inches(0.55), light=True)

    # "Prepared for" badge — only when a hospital name is provided.
    # Sits between the logo and the eyebrow so it's the first thing the
    # reader's eye lands on after the brand.
    eyebrow_y = Inches(2.4)
    if hospital_name and prepared_for_line:
        badge_w = Inches(8.0)
        badge_h = Inches(0.55)
        badge = add_rounded(slide, Inches(0.6), Inches(1.6), badge_w, badge_h, GREEN)
        badge.adjustments[0] = 0.5
        add_text(slide, Inches(0.6), Inches(1.6), badge_w, badge_h,
                 prepared_for_line.upper(),
                 size=13, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Eyebrow
    add_text(slide, Inches(0.6), eyebrow_y, Inches(8), Inches(0.5),
             "HOSPITAL CASE-MEETING PLATFORM",
             size=14, color=GREEN, bold=True)

    # Main title (massive, multi-line)
    add_text(slide, Inches(0.6), Inches(2.95), Inches(10), Inches(2.5),
             "Stop scheduling.\nStart deciding.",
             size=60, color=WHITE, bold=True)

    # Subhead — softly mentions the hospital if provided
    if hospital_name:
        subhead = (
            f"A walkthrough for {hospital_name} — how BioMedMeet turns weekly "
            "multidisciplinary case meetings into structured, audit-ready "
            "decisions without the calendaring, minute-taking and follow-up "
            "overhead."
        )
    else:
        subhead = (
            "BioMedMeet turns weekly multidisciplinary case meetings into "
            "structured, audit-ready decisions — without the calendaring, "
            "minute-taking and follow-up overhead."
        )
    add_text(slide, Inches(0.6), Inches(5.0), Inches(9), Inches(1.2),
             subhead, size=18, color=LIGHT)

    # CTA pill
    cta = add_rounded(slide, Inches(0.6), Inches(6.3), Inches(3.3), Inches(0.6),
                      GREEN)
    cta.adjustments[0] = 0.5
    add_text(slide, Inches(0.6), Inches(6.3), Inches(3.3), Inches(0.6),
             "Request your 15-minute demo",
             size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Footer ribbon (mentions the hospital when set)
    footer_left = (
        f"Demo deck  ·  Prepared for {hospital_name}"
        if hospital_name
        else "Demo deck  ·  Prepared for hospital leadership"
    )
    add_text(slide, Inches(0.6), Inches(7.05), Inches(12), Inches(0.3),
             f"{footer_left}  ·  Demo@BioMedMeet.com  ·  biomedmeet.com",
             size=10, color=LIGHT)


def make_overview(prs, *, hospital_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)

    # Top accent bar
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    # Section eyebrow + title
    eyebrow = (
        f"WHAT BIOMEDMEET MEANS FOR {hospital_name.upper()}"
        if hospital_name else "WHAT IS BIOMEDMEET?"
    )
    add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             eyebrow, size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(11), Inches(1.0),
             "The case-meeting platform built for hospital teams.",
             size=34, color=NAVY, bold=True)

    # 2-column layout: narrative left, KPI cards right
    add_text(slide, Inches(0.6), Inches(2.95), Inches(6.7), Inches(3.5),
             "Tumor boards, MDTs, M&Ms, departmental rounds — every hospital "
             "runs them, and every hospital wastes 15+ hours per team per week "
             "around them.\n\n"
             "BioMedMeet replaces the spreadsheets, email threads, ad-hoc "
             "Teams invites and post-meeting follow-ups with one purpose-built "
             "workflow: schedule, present cases, log decisions, generate a "
             "signed PDF summary, distribute it automatically.\n\n"
             "Self-hosted on your own network. No PHI ever leaves the hospital. "
             "Microsoft Teams integration. Holiday-aware recurring schedules. "
             "RSVP via email or in-app. Decision-log audit trail.",
             size=14, color=SLATE)

    # Right column — 3 KPI cards
    kpis = [
        ("18 hrs",    "saved per team per week",     GREEN,  TINT_GREEN),
        ("$187,200",  "clinician time recovered / yr", OLIVE,  TINT_OLIVE),
        ("~1.6 mo",   "payback period vs. licence",  PURPLE, TINT_PURPLE),
    ]
    x0 = Inches(8.0)
    for i, (stat, label, accent, tint) in enumerate(kpis):
        y = Inches(2.95 + i * 1.35)
        add_rounded(slide, x0, y, Inches(4.6), Inches(1.15), tint)
        # accent stripe
        add_rect(slide, x0, y, Inches(0.10), Inches(1.15), accent)
        add_text(slide, x0 + Inches(0.30), y + Inches(0.12),
                 Inches(4.2), Inches(0.55),
                 stat, size=30, color=accent, bold=True)
        add_text(slide, x0 + Inches(0.30), y + Inches(0.68),
                 Inches(4.2), Inches(0.4),
                 label, size=12, color=SLATE)

    add_footer(slide, 2, 10)


def make_problem_solution(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, IVORY)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    add_text(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "WHY HOSPITALS BUY IT", size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "From scattered tools  →  one structured workflow.",
             size=32, color=NAVY, bold=True)

    # Two big columns: BEFORE  →  AFTER
    col_w = Inches(5.7)
    col_h = Inches(4.5)
    add_rounded(slide, Inches(0.6), Inches(2.85), col_w, col_h, WHITE,
                line_color=LIGHT)
    add_rounded(slide, Inches(6.95), Inches(2.85), col_w, col_h, NAVY)

    # BEFORE
    add_text(slide, Inches(0.95), Inches(3.05), col_w - Inches(0.7), Inches(0.5),
             "BEFORE", size=11, color=OLIVE, bold=True)
    add_text(slide, Inches(0.95), Inches(3.35), col_w - Inches(0.7), Inches(0.6),
             "Today, every Monday morning…",
             size=20, color=NAVY, bold=True)
    before_items = [
        "5 different Outlook threads scheduling one MDT",
        "Agenda Word doc pasted into a chat 10 minutes before",
        "Decisions captured in a scribe's notebook",
        "Follow-ups manually dispatched to oncology, radiology, surgery",
        "No audit trail when a decision is later challenged",
    ]
    for i, t in enumerate(before_items):
        y = Inches(4.1 + i * 0.45)
        add_text(slide, Inches(1.05), y, Inches(0.25), Inches(0.35),
                 "✕", size=14, color=OLIVE, bold=True)
        add_text(slide, Inches(1.35), y, col_w - Inches(1.0), Inches(0.4),
                 t, size=13, color=SLATE)

    # AFTER
    add_text(slide, Inches(7.30), Inches(3.05), col_w - Inches(0.7), Inches(0.5),
             "AFTER", size=11, color=GREEN, bold=True)
    add_text(slide, Inches(7.30), Inches(3.35), col_w - Inches(0.7), Inches(0.6),
             "With BioMedMeet…",
             size=20, color=WHITE, bold=True)
    after_items = [
        "Recurring MDT scheduled once, Teams link auto-attached",
        "Agenda pulled from the patient list, in two clicks",
        "Decisions logged live, signed off by the chair",
        "PDF summary auto-mailed to attendees & EMR",
        "Every decision time-stamped, attributable, exportable",
    ]
    for i, t in enumerate(after_items):
        y = Inches(4.1 + i * 0.45)
        add_text(slide, Inches(7.40), y, Inches(0.25), Inches(0.35),
                 "✓", size=14, color=GREEN, bold=True)
        add_text(slide, Inches(7.70), y, col_w - Inches(1.0), Inches(0.4),
                 t, size=13, color=LIGHT)

    add_footer(slide, 3, 10)


def make_feature_grid(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    add_text(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "PLATFORM CAPABILITIES", size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "Six pillars. One platform. Zero spreadsheets.",
             size=30, color=NAVY, bold=True)

    features = [
        ("01", "Smart scheduling",
         "Recurring meetings (daily / weekly / monthly-on), holiday-aware "
         "for US, India and UK calendars, with Teams links generated on save.",
         GREEN, TINT_GREEN),
        ("02", "Patient case rooms",
         "Patient cases attached to a meeting with diagnoses, MRN, allergies, "
         "medications, files and a 5-tab case room for the discussion.",
         OLIVE, TINT_OLIVE),
        ("03", "Decision log",
         "Every decision captured with the responsible doctor, priority, "
         "follow-up date and free-text rationale — searchable forever.",
         PURPLE, TINT_PURPLE),
        ("04", "Treatment plans",
         "A single source of truth per patient. Latest plan rolls up to the "
         "patient detail page; older plans become an audit history.",
         NAVY, TINT_NAVY),
        ("05", "PDF summaries",
         "One-click PDF summary at meeting close — attendees, decisions, "
         "treatment plans, agenda items. Mailed automatically.",
         GREEN, TINT_GREEN),
        ("06", "Inbound RSVP",
         "Participants RSVP from email links *and* by replying to the invite. "
         "An admin RSVP log surfaces every parsed reply for audit.",
         OLIVE, TINT_OLIVE),
    ]
    card_w, card_h = Inches(4.15), Inches(1.85)
    x_start, y_start = Inches(0.6), Inches(2.85)
    gap_x, gap_y = Inches(0.10), Inches(0.18)

    for i, (num, title, desc, accent, tint) in enumerate(features):
        row, col = divmod(i, 3)
        x = x_start + (card_w + gap_x) * col
        y = y_start + (card_h + gap_y) * row
        add_rounded(slide, x, y, card_w, card_h, tint)
        add_rect(slide, x, y, Inches(0.10), card_h, accent)
        add_text(slide, x + Inches(0.30), y + Inches(0.12),
                 Inches(1.5), Inches(0.4),
                 num, size=22, color=accent, bold=True)
        add_text(slide, x + Inches(0.95), y + Inches(0.15),
                 card_w - Inches(1.0), Inches(0.4),
                 title, size=15, color=NAVY, bold=True)
        add_text(slide, x + Inches(0.30), y + Inches(0.62),
                 card_w - Inches(0.5), card_h - Inches(0.7),
                 desc, size=11, color=SLATE)

    add_footer(slide, 4, 10)


def make_feature_deep_workflow(prs, step_thumbs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, IVORY)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    add_text(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "THE 4-STEP MEETING WORKFLOW", size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "Schedule  →  Discuss  →  Decide  →  Distribute.",
             size=30, color=NAVY, bold=True)

    steps = [
        ("Schedule", "Recurring meetings with Teams links auto-attached, "
                     "holiday-aware date validation, and email invites with iCal.",
                     GREEN),
        ("Discuss",  "Live 5-tab case room: Patients, Agenda, Files, "
                     "Treatment Plans, Decisions. Built for an MDT panel.",
                     OLIVE),
        ("Decide",   "Decision log captures the chair-approved outcome, "
                     "owner, priority and follow-up — every entry signed.",
                     PURPLE),
        ("Distribute","Auto PDF + auto-email to attendees the moment the "
                      "meeting is closed. Searchable history per patient.",
                      NAVY),
    ]
    cell_w = Inches(3.0)
    cell_h = Inches(3.8)
    x0 = Inches(0.6)
    y0 = Inches(2.85)
    gap = Inches(0.10)

    for i, (title, desc, accent) in enumerate(steps):
        x = x0 + (cell_w + gap) * i
        add_rounded(slide, x, y0, cell_w, cell_h, WHITE, line_color=LIGHT)
        # Numbered chip
        chip = add_rounded(slide, x + Inches(0.30), y0 + Inches(0.25),
                           Inches(0.6), Inches(0.6), accent)
        chip.adjustments[0] = 0.5
        add_text(slide, x + Inches(0.30), y0 + Inches(0.25),
                 Inches(0.6), Inches(0.6),
                 str(i + 1), size=22, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(1.05), y0 + Inches(0.28),
                 cell_w - Inches(1.1), Inches(0.6),
                 title, size=20, color=NAVY, bold=True)
        add_text(slide, x + Inches(0.30), y0 + Inches(1.05),
                 cell_w - Inches(0.5), Inches(1.4),
                 desc, size=11, color=SLATE)

        # Step thumbnail (from public/marketing/thumbs/stepN.jpg)
        thumb = step_thumbs[i]
        if thumb.exists():
            slide.shapes.add_picture(str(thumb),
                                     x + Inches(0.30),
                                     y0 + Inches(2.40),
                                     width=cell_w - Inches(0.5),
                                     height=Inches(1.25))

    add_footer(slide, 5, 10)


def make_security_compliance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    add_text(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "ARCHITECTURE & COMPLIANCE", size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "Designed to satisfy your CISO before your CMO.",
             size=30, color=NAVY, bold=True)

    pillars = [
        ("On-prem by default",
         "Self-hosted Docker stack on your hospital network. "
         "PHI never leaves the building.", GREEN),
        ("Microsoft 365 native",
         "Graph API for Teams meetings, SMTP for invites, "
         "iCal for two-way calendar sync.", OLIVE),
        ("Audit-ready",
         "Every decision timestamped, attributable, exportable. "
         "Soft-deletes preserve history.", PURPLE),
        ("Role-based access",
         "Organizer, doctor, nurse, admin. Field-level guardrails on "
         "patient updates, agenda changes & decisions.", NAVY),
    ]
    cw, ch = Inches(5.8), Inches(1.85)
    for i, (title, desc, accent) in enumerate(pillars):
        r, c = divmod(i, 2)
        x = Inches(0.6) + (cw + Inches(0.2)) * c
        y = Inches(2.85) + (ch + Inches(0.25)) * r
        add_rounded(slide, x, y, cw, ch, WHITE, line_color=LIGHT)
        add_rect(slide, x, y, Inches(0.10), ch, accent)
        add_text(slide, x + Inches(0.35), y + Inches(0.20),
                 cw - Inches(0.5), Inches(0.5),
                 title, size=18, color=NAVY, bold=True)
        add_text(slide, x + Inches(0.35), y + Inches(0.78),
                 cw - Inches(0.5), ch - Inches(0.9),
                 desc, size=12, color=SLATE)

    add_footer(slide, 6, 10)


def make_pricing(prs, *, hospital_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    eyebrow = (
        f"PRICING FOR {hospital_name.upper()}" if hospital_name else "PRICING"
    )
    add_text(slide, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
             eyebrow, size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "Own it. Don't rent it.",
             size=34, color=NAVY, bold=True)
    add_text(slide, Inches(0.6), Inches(2.35), Inches(12), Inches(0.5),
             "One-time perpetual licence. Self-hosted. No per-user fees. "
             "No monthly bill. Two-month warranty trial.",
             size=14, color=SLATE)

    # The big price card
    card = add_rounded(slide, Inches(0.6), Inches(3.05),
                       Inches(6.4), Inches(3.85), WHITE, line_color=LIGHT)
    add_text(slide, Inches(0.95), Inches(3.30), Inches(5.5), Inches(0.4),
             "PERPETUAL LICENCE  ·  SELF-HOSTED",
             size=10, color=GREEN, bold=True)
    pricing_title = (
        f"BioMedMeet for {hospital_name}" if hospital_name
        else "BioMedMeet for hospitals"
    )
    add_text(slide, Inches(0.95), Inches(3.65), Inches(5.5), Inches(0.6),
             pricing_title,
             size=16, color=SLATE, bold=True)
    add_text(slide, Inches(0.95), Inches(4.15), Inches(5.5), Inches(1.4),
             "$25,000",
             size=72, color=NAVY, bold=True)
    add_text(slide, Inches(0.95), Inches(5.55), Inches(5.5), Inches(0.4),
             "+ applicable taxes  ·  one-time  ·  no recurring fees",
             size=11, color=SLATE, italic=True)
    bullets = [
        "Unlimited users · unlimited meetings · unlimited patients",
        "Microsoft Teams + holiday-aware recurring calendar",
        "Decision log · treatment plans · PDF summaries · audit trail",
        "Two-month warranty trial — full refund window",
    ]
    for i, b in enumerate(bullets):
        y = Inches(5.95 + i * 0.24)
        add_text(slide, Inches(0.95), y, Inches(0.2), Inches(0.25),
                 "✓", size=11, color=GREEN, bold=True)
        add_text(slide, Inches(1.15), y, Inches(5.4), Inches(0.25),
                 b, size=11, color=SLATE)

    # Right: "what's not in the price"
    right = add_rounded(slide, Inches(7.30), Inches(3.05),
                        Inches(5.4), Inches(3.85), NAVY)
    add_text(slide, Inches(7.60), Inches(3.30), Inches(4.8), Inches(0.4),
             "WHAT YOU DON'T PAY",
             size=10, color=GREEN, bold=True)
    add_text(slide, Inches(7.60), Inches(3.65), Inches(4.8), Inches(0.6),
             "Compare against the SaaS alternative.",
             size=18, color=WHITE, bold=True)
    saas_rows = [
        ("Per-user monthly fees",       "$0"),
        ("Per-meeting / per-case fees", "$0"),
        ("Cloud storage add-ons",       "$0"),
        ("Annual renewal markup",       "$0"),
        ("Data-egress charges",         "$0"),
        ("Lock-in clause",              "None"),
    ]
    for i, (k, v) in enumerate(saas_rows):
        y = Inches(4.45 + i * 0.38)
        add_text(slide, Inches(7.60), y, Inches(3.5), Inches(0.35),
                 k, size=12, color=LIGHT)
        add_text(slide, Inches(11.20), y, Inches(1.3), Inches(0.35),
                 v, size=12, color=GREEN, bold=True, align=PP_ALIGN.RIGHT)

    add_footer(slide, 7, 10)


def make_roi(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, IVORY)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    add_text(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "RETURN ON INVESTMENT", size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "Pays for itself before the warranty trial ends.",
             size=30, color=NAVY, bold=True)

    # Three big stat tiles
    tiles = [
        ("18 hrs / week",      "Hours reclaimed per MDT", GREEN,  TINT_GREEN),
        ("$187,200 / year",    "Clinician time recovered", OLIVE,  TINT_OLIVE),
        ("1.6 months",         "Payback period",          PURPLE, TINT_PURPLE),
    ]
    tw, th = Inches(4.0), Inches(1.85)
    for i, (stat, label, accent, tint) in enumerate(tiles):
        x = Inches(0.6) + (tw + Inches(0.2)) * i
        y = Inches(2.85)
        add_rounded(slide, x, y, tw, th, tint)
        add_rect(slide, x, y, Inches(0.10), th, accent)
        add_text(slide, x + Inches(0.30), y + Inches(0.15),
                 tw - Inches(0.5), Inches(0.7),
                 stat, size=28, color=accent, bold=True)
        add_text(slide, x + Inches(0.30), y + Inches(0.95),
                 tw - Inches(0.5), Inches(0.4),
                 label, size=12, color=SLATE)

    # "Where the 18 hours go" bar chart
    add_text(slide, Inches(0.6), Inches(5.10), Inches(8), Inches(0.4),
             "Where the 18 hours go each week",
             size=18, color=NAVY, bold=True)
    add_text(slide, Inches(0.6), Inches(5.45), Inches(11), Inches(0.3),
             "Time-and-motion estimates from clinical teams running 1–2 case conferences/week.",
             size=10, color=SLATE, italic=True)

    bars = [
        ("Scheduling & calendaring",      5.0, GREEN),
        ("Agenda prep & case pulls",      4.0, OLIVE),
        ("Minute-taking & decision log",  3.5, PURPLE),
        ("RSVP chase & no-show recovery", 3.0, NAVY),
        ("Post-meeting follow-up emails", 2.5, SLATE),
    ]
    max_hours = 5.0
    bar_x0 = Inches(3.5)
    bar_w_max = Inches(6.5)
    for i, (label, hrs, accent) in enumerate(bars):
        y = Inches(5.85 + i * 0.26)
        add_text(slide, Inches(0.6), y, Inches(2.8), Inches(0.25),
                 label, size=11, color=NAVY)
        # bar track
        add_rounded(slide, bar_x0, y + Inches(0.04), bar_w_max, Inches(0.18),
                    LIGHT)
        # bar fill
        fw = Inches(6.5 * hrs / max_hours)
        bar = add_rounded(slide, bar_x0, y + Inches(0.04), fw, Inches(0.18),
                          accent)
        bar.adjustments[0] = 0.4
        # value label at right
        add_text(slide, bar_x0 + bar_w_max + Inches(0.1), y,
                 Inches(1.0), Inches(0.25),
                 f"{hrs:.1f} hrs", size=11, color=NAVY, bold=True)

    add_footer(slide, 8, 10)


def make_revenue_business(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, CREAM)
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.18), NAVY)
    draw_logo(slide, Inches(0.5), Inches(0.4), height=Inches(0.5))

    add_text(slide, Inches(0.6), Inches(1.2), Inches(8), Inches(0.4),
             "REVENUE & THROUGHPUT BENEFITS", size=12, color=GREEN, bold=True)
    add_text(slide, Inches(0.6), Inches(1.55), Inches(12), Inches(1.0),
             "Beyond hours — what BioMedMeet unlocks.",
             size=30, color=NAVY, bold=True)

    benefits = [
        ("Faster decisions",
         "Decision-to-distribution time falls from days to minutes. "
         "Patients reach treatment sooner; throughput per MDT rises.",
         "+1 case / MDT / week", GREEN, TINT_GREEN),
        ("Higher-value MDTs",
         "Structured case rooms let panels handle 30–40% more cases in "
         "the same calendar slot. No time lost to scheduling friction.",
         "30-40% capacity", OLIVE, TINT_OLIVE),
        ("Stronger billing & coding",
         "Auto-generated, signed PDF summaries feed coders with structured "
         "decisions — fewer rejected claims, cleaner CPT documentation.",
         "Lower denials", PURPLE, TINT_PURPLE),
        ("Audit & accreditation",
         "Joint Commission / NABH-ready decision trail. One click to export "
         "any meeting's full history for review.",
         "Audit-ready", NAVY, TINT_NAVY),
    ]
    cw, ch = Inches(5.8), Inches(2.25)
    for i, (title, desc, chip, accent, tint) in enumerate(benefits):
        r, c = divmod(i, 2)
        x = Inches(0.6) + (cw + Inches(0.2)) * c
        y = Inches(2.85) + (ch + Inches(0.2)) * r
        add_rounded(slide, x, y, cw, ch, tint)
        add_rect(slide, x, y, Inches(0.10), ch, accent)

        add_text(slide, x + Inches(0.35), y + Inches(0.18),
                 cw - Inches(2.4), Inches(0.5),
                 title, size=17, color=NAVY, bold=True)

        # Chip
        chip_w = Inches(1.85)
        chip_shape = add_rounded(slide,
                                 x + cw - chip_w - Inches(0.25),
                                 y + Inches(0.22),
                                 chip_w, Inches(0.35),
                                 accent)
        chip_shape.adjustments[0] = 0.5
        add_text(slide, x + cw - chip_w - Inches(0.25), y + Inches(0.22),
                 chip_w, Inches(0.35), chip,
                 size=10, color=WHITE, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        add_text(slide, x + Inches(0.35), y + Inches(0.80),
                 cw - Inches(0.6), ch - Inches(0.9),
                 desc, size=12, color=SLATE)

    add_footer(slide, 9, 10)


def make_cta(prs, *, hospital_name=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, prs.slide_width, prs.slide_height, NAVY)

    # Decorative slab
    slab = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                  Inches(-2), Inches(5.5),
                                  Inches(20), Inches(3.5))
    slab.rotation = -6
    slab.fill.solid()
    slab.fill.fore_color.rgb = GREEN
    slab.line.fill.background()
    slab.shadow.inherit = False

    draw_logo(slide, Inches(0.6), Inches(0.5), height=Inches(0.55), light=True)

    add_text(slide, Inches(0.6), Inches(2.0), Inches(11), Inches(0.5),
             "READY TO SEE IT LIVE?", size=14, color=GREEN, bold=True)
    cta_title = (
        f"Let's book a 15-minute walkthrough\nfor {hospital_name}."
        if hospital_name
        else "Let's book a 15-minute walkthrough."
    )
    add_text(slide, Inches(0.6), Inches(2.55), Inches(11), Inches(2.0),
             cta_title, size=44 if hospital_name else 48,
             color=WHITE, bold=True)
    cta_sub = (
        f"We'll bring a working sandbox so {hospital_name}'s team can click "
        "around — schedule a recurring tumor board, run a sample case, "
        "generate the PDF — all in real time."
        if hospital_name
        else
        "We'll bring a working sandbox so your team can click around — "
        "schedule a recurring tumor board, run a sample case, generate "
        "the PDF — all in real time."
    )
    add_text(slide, Inches(0.6), Inches(4.20), Inches(11), Inches(1.2),
             cta_sub, size=16, color=LIGHT)

    # CTA buttons row
    cta1 = add_rounded(slide, Inches(0.6), Inches(6.0),
                       Inches(4.2), Inches(0.7), WHITE)
    cta1.adjustments[0] = 0.5
    add_text(slide, Inches(0.6), Inches(6.0), Inches(4.2), Inches(0.7),
             "Email  Demo@BioMedMeet.com",
             size=14, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cta2 = add_rounded(slide, Inches(5.0), Inches(6.0),
                       Inches(4.2), Inches(0.7), OLIVE)
    cta2.adjustments[0] = 0.5
    add_text(slide, Inches(5.0), Inches(6.0), Inches(4.2), Inches(0.7),
             "Visit  biomedmeet.com",
             size=14, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(0.6), Inches(7.10), Inches(12), Inches(0.3),
             "Thank you  ·  BioMedMeet  ·  biomedmeet.com",
             size=10, color=LIGHT)


# ----------------------------------------------------------------------------- 
# Main
# ----------------------------------------------------------------------------- 

def _slug(name: str) -> str:
    """Turn 'Ochsner Health' → 'Ochsner_Health' for filenames."""
    s = re.sub(r"[^A-Za-z0-9]+", "_", name).strip("_")
    return s or "Hospital"


def parse_args():
    p = argparse.ArgumentParser(
        description="Generate the BioMedMeet demo PowerPoint deck.")
    p.add_argument("--hospital-name", default=None,
                   help="Personalises the cover, overview, pricing & CTA "
                        "(e.g. \"Ochsner Health\").")
    p.add_argument("--month", default=None,
                   help="Month name shown on the 'Prepared for' badge "
                        "(default: current month).")
    p.add_argument("--year", type=int, default=None,
                   help="Year shown on the 'Prepared for' badge "
                        "(default: current year).")
    p.add_argument("--output", default=None,
                   help="Output .pptx path. Defaults to "
                        "/app/docs/BioMedMeet_Demo_Deck[_<Hospital>].pptx")
    return p.parse_args()


def main():
    args = parse_args()
    now = datetime.now()
    month = args.month or now.strftime("%B")
    year = args.year or now.year
    hospital = args.hospital_name.strip() if args.hospital_name else None
    prepared_for = (
        f"Prepared for {hospital}  ·  {month} {year}" if hospital else None
    )

    if args.output:
        out_path = Path(args.output)
    elif hospital:
        out_path = OUT_PATH.with_name(
            f"BioMedMeet_Demo_Deck_{_slug(hospital)}.pptx"
        )
    else:
        out_path = OUT_PATH

    prs = Presentation()
    prs.slide_width = Inches(13.333)   # 16:9
    prs.slide_height = Inches(7.5)

    step_thumb_dir = Path("/app/frontend/public/marketing/thumbs")
    step_thumbs = [step_thumb_dir / f"step{i}.jpg" for i in range(1, 5)]

    make_cover(prs, hospital_name=hospital,
               prepared_for_line=prepared_for)             # 1
    make_overview(prs, hospital_name=hospital)             # 2
    make_problem_solution(prs)                             # 3
    make_feature_grid(prs)                                 # 4
    make_feature_deep_workflow(prs, step_thumbs)           # 5
    make_security_compliance(prs)                          # 6
    make_pricing(prs, hospital_name=hospital)              # 7
    make_roi(prs)                                          # 8
    make_revenue_business(prs)                             # 9
    make_cta(prs, hospital_name=hospital)                  # 10

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    print(f"Saved: {out_path}  ({len(prs.slides)} slides)")
    if hospital:
        print(f"Personalised for: {hospital}  ·  {month} {year}")


if __name__ == "__main__":
    main()
